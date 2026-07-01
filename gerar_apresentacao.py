"""
Gera uma apresentação PowerPoint do TCC.

O arquivo é opcional para executar o sistema principal. Ele existe para facilitar
a criação de uma apresentação com os pontos centrais do projeto.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def adicionar_caixa_texto(slide, esquerda, topo, largura, altura, texto, tamanho=20, cor=None, negrito=False):
    """Cria uma caixa de texto simples e retorna o parágrafo criado."""

    caixa = slide.shapes.add_textbox(Inches(esquerda), Inches(topo), Inches(largura), Inches(altura))
    paragrafo = caixa.text_frame.paragraphs[0]
    paragrafo.text = texto
    paragrafo.font.name = "Arial"
    paragrafo.font.size = Pt(tamanho)
    paragrafo.font.bold = negrito
    if cor:
        paragrafo.font.color.rgb = cor
    return paragrafo


def adicionar_slide_topicos(apresentacao, titulo, topicos, cores):
    """Adiciona um slide com título e lista de tópicos."""

    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    fundo = slide.background.fill
    fundo.solid()
    fundo.fore_color.rgb = cores["fundo"]

    adicionar_caixa_texto(slide, 0.8, 0.5, 11.8, 0.7, titulo, 30, cores["titulo"], True)

    caixa = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.2), Inches(4.8))
    quadro = caixa.text_frame
    quadro.word_wrap = True

    for indice, topico in enumerate(topicos):
        paragrafo = quadro.paragraphs[0] if indice == 0 else quadro.add_paragraph()
        paragrafo.text = f"• {topico}"
        paragrafo.font.name = "Arial"
        paragrafo.font.size = Pt(18)
        paragrafo.font.color.rgb = cores["texto"]
        paragrafo.space_before = Pt(10)

    return slide


def gerar_apresentacao():
    """Monta e salva o arquivo .pptx da apresentação."""

    apresentacao = Presentation()
    apresentacao.slide_width = Inches(13.333)
    apresentacao.slide_height = Inches(7.5)

    cores = {
        "fundo": RGBColor(255, 255, 255),
        "titulo": RGBColor(242, 143, 13),
        "texto": RGBColor(40, 40, 40),
        "roxo": RGBColor(90, 45, 130),
        "cinza": RGBColor(120, 120, 120),
    }

    # Slide de capa.
    capa = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    capa.background.fill.solid()
    capa.background.fill.fore_color.rgb = cores["fundo"]

    titulo = adicionar_caixa_texto(
        capa,
        2.2,
        2.0,
        9.0,
        0.8,
        "Servidor de Impressão Automatizada",
        34,
        cores["titulo"],
        True,
    )
    titulo.alignment = PP_ALIGN.CENTER

    subtitulo = adicionar_caixa_texto(
        capa,
        2.2,
        3.0,
        9.0,
        0.6,
        "Revitalização de impressoras legadas com interface web e pagamento Pix",
        18,
        cores["cinza"],
    )
    subtitulo.alignment = PP_ALIGN.CENTER

    autor = adicionar_caixa_texto(capa, 2.2, 4.3, 9.0, 0.6, "Autor: Lucas Bastos Rezende", 16, cores["texto"], True)
    autor.alignment = PP_ALIGN.CENTER

    adicionar_slide_topicos(
        apresentacao,
        "Problema",
        [
            "Impressoras antigas ainda funcionam, mas muitas são descartadas por falta de conectividade.",
            "O envio manual de arquivos por pendrive, e-mail ou mensagens gera demora e risco de segurança.",
            "A cobrança manual pode causar erro de valor e dificuldade de controle.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Solução proposta",
        [
            "Criar um servidor local para receber arquivos pelo navegador.",
            "Permitir configuração de páginas, cópias, orientação, cor e tamanho de papel.",
            "Liberar a impressão depois do pagamento via Pix ou confirmação manual.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Arquitetura",
        [
            "Interface em HTML, CSS e JavaScript.",
            "Servidor Python com Flask.",
            "Impressão silenciosa com SumatraPDF no Windows.",
            "Cobrança Pix pelo Mercado Pago ou QR Code estático.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Funcionalidades",
        [
            "Envio de PDF, documentos de texto e imagens.",
            "Pré-visualização de PDFs e imagens.",
            "Seleção de páginas e quantidade de cópias.",
            "Modo de demonstração para apresentação do TCC.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Segurança e privacidade",
        [
            "O arquivo enviado fica salvo somente de forma temporária.",
            "Depois da tentativa de impressão, o arquivo é apagado automaticamente.",
            "O sistema evita circulação de documentos por pendrive ou aplicativos de mensagem.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Trabalhos futuros",
        [
            "Painel administrativo com histórico de impressões.",
            "Fila com múltiplas impressoras.",
            "Webhook do Mercado Pago para substituir a consulta por intervalo.",
            "Versão Linux usando CUPS.",
        ],
        cores,
    )

    adicionar_slide_topicos(
        apresentacao,
        "Conclusão",
        [
            "O projeto demonstra que software pode modernizar impressoras comuns.",
            "A solução reduz trabalho manual e organiza o atendimento.",
            "O sistema combina reaproveitamento de hardware, interface web e pagamento digital.",
        ],
        cores,
    )

    caminho_saida = "Apresentacao_TCC_Unisapiens.pptx"
    apresentacao.save(caminho_saida)
    print(f"Apresentação salva em: {os.path.abspath(caminho_saida)}")


if __name__ == "__main__":
    gerar_apresentacao()
